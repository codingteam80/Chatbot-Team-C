"""Answer action buttons: copy, regenerate, and export."""

import csv
import html
import json
import re
import time
from datetime import datetime
from io import BytesIO, StringIO
from pathlib import Path

import streamlit as st
import streamlit.components.v1 as components

try:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches as DocxInches
    from docx.shared import Pt as DocxPt
except Exception:
    Document = None
    WD_ALIGN_PARAGRAPH = None
    DocxInches = None
    DocxPt = None

try:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
except Exception:
    Workbook = None
    Alignment = None
    Border = None
    Font = None
    PatternFill = None
    Side = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.util import Pt as PptxPt
except Exception:
    Presentation = None
    PptxInches = None
    PptxPt = None

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
except Exception:
    colors = None
    letter = None
    ParagraphStyle = None
    getSampleStyleSheet = None
    Paragraph = None
    SimpleDocTemplate = None
    Spacer = None

from ui.source_ui import (
    clean_preview_text,
    deduplicate_sources,
    make_static_file_link,
    normalize_source_item,
)

SIDEBAR_SOURCES_KEY = "sidebar_sources"
SOURCE_PREVIEW_LIMIT = 360
EXPORT_TIMESTAMP_FORMAT = "%Y%m%d_%H%M%S"
PENDING_EXPORT_REQUEST_KEY = "pending_export_request"
READY_EXPORT_FILE_KEY = "ready_export_file"
EXPORT_FRONTEND_FLUSH_SECONDS = 0.12
ACTION_DEBUG_FILE = Path(__file__).resolve().parents[1] / "reports" / "ui_rag_debug.txt"


def append_action_debug(title, lines=None):
    # Confirms whether action_ui.py receives the button click.
    try:
        ACTION_DEBUG_FILE.parent.mkdir(parents=True, exist_ok=True)
        log_lines = [
            "",
            "=" * 80,
            str(title or "ACTION DEBUG"),
            "=" * 80,
            f"Time: {time.strftime('%Y-%m-%d %H:%M:%S')}",
        ]

        for line in lines or []:
            log_lines.append(str(line))

        with ACTION_DEBUG_FILE.open("a", encoding="utf-8") as file:
            file.write("\n".join(log_lines))
            file.write("\n")

        print(f"ACTION DEBUG: {title} -> {ACTION_DEBUG_FILE.resolve()}", flush=True)
    except Exception as error:
        print(f"ACTION DEBUG WRITE FAILED: {error}", flush=True)


print(f"ACTIVE ACTION_UI FILE: {Path(__file__).resolve()}", flush=True)
print(f"ACTION DEBUG FILE TARGET: {ACTION_DEBUG_FILE.resolve()}", flush=True)

EXPORT_TYPES = {
    "PDF": ("pdf", "application/pdf"),
    "DOCX": ("docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "TXT": ("txt", "text/plain"),
    "Markdown": ("md", "text/markdown"),
    "CSV": ("csv", "text/csv"),
    "XLSX": ("xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    "XLS": ("xls", "application/vnd.ms-excel"),
    "PPTX": ("pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "PPT": ("ppt", "application/vnd.ms-powerpoint"),
}


def safe_html(value):
    # Escape text for custom HTML blocks.
    return html.escape(str(value or ""), quote=True)


def copy_button(text, key):
    # Browser clipboard needs a tiny component iframe; main CSS cannot style inside iframes.
    safe_text = json.dumps(text or "")
    safe_key = safe_html(key)

    components.html(
        f"""
<!DOCTYPE html>
<html>
<head>
<style>
html,body{{margin:0;padding:0;width:34px;height:34px;overflow:hidden;background:transparent;}}
button{{width:34px;height:34px;border:0;border-radius:10px;background:transparent;color:rgba(255,255,255,.82);font-size:15px;cursor:pointer;}}
button:hover{{background:rgba(255,255,255,.08);color:#fff;}}
</style>
</head>
<body>
<button id="copy-btn-{safe_key}" title="Copy" onclick="copyText()">⧉</button>
<script>
function fallbackCopy(text){{
  const el=document.createElement('textarea');
  el.value=text; el.style.position='fixed'; el.style.opacity='0';
  document.body.appendChild(el); el.focus(); el.select();
  try{{document.execCommand('copy');}}catch(err){{}}
  document.body.removeChild(el);
}}
function showCopied(){{
  const btn=document.getElementById('copy-btn-{safe_key}');
  btn.innerText='✓'; setTimeout(function(){{btn.innerText='⧉';}},900);
}}
function copyText(){{
  const text={safe_text};
  if(navigator.clipboard && window.isSecureContext){{
    navigator.clipboard.writeText(text).then(showCopied).catch(function(){{fallbackCopy(text);showCopied();}});
  }} else {{ fallbackCopy(text); showCopied(); }}
}}
</script>
</body>
</html>
""",
        height=34,
        width=34,
        scrolling=False,
    )


def get_export_timestamp():
    # One timestamp per click/rerun for readable filenames and export metadata.
    return datetime.now().strftime(EXPORT_TIMESTAMP_FORMAT)


def title_case_words(text):
    # Make a readable title without depending on the LLM.
    small_words = {"a", "an", "and", "as", "at", "by", "for", "from", "in", "of", "on", "or", "the", "to", "with"}
    words = str(text or "").split()
    result = []

    for index, word in enumerate(words):
        if word.isupper() and len(word) <= 6:
            result.append(word)
            continue

        lower_word = word.lower()

        if index > 0 and lower_word in small_words:
            result.append(lower_word)
        else:
            result.append(lower_word[:1].upper() + lower_word[1:])

    return " ".join(result).strip()


def make_response_title(question=None):
    # Create a short topic title related to the question, but not the exact question.
    text = str(question or "").strip()

    if not text:
        return "InknowVa Response"

    text = re.sub(r"[?!.]+$", "", text)
    text = re.sub(r"\s+", " ", text).strip()

    replacements = [
        (r"^(please\s+)?(can|could|would)\s+you\s+", ""),
        (r"^(please\s+)?tell\s+me\s+about\s+", ""),
        (r"^(please\s+)?explain\s+", ""),
        (r"^what\s+is\s+", ""),
        (r"^what\s+are\s+", ""),
        (r"^who\s+is\s+", ""),
        (r"^who\s+are\s+", ""),
        (r"^when\s+did\s+", ""),
        (r"^when\s+is\s+", ""),
        (r"^where\s+is\s+", ""),
        (r"^where\s+are\s+", ""),
        (r"^why\s+is\s+", ""),
        (r"^why\s+are\s+", ""),
        (r"^how\s+does\s+", ""),
        (r"^how\s+do\s+", ""),
        (r"^how\s+did\s+", ""),
        (r"^ano\s+ang\s+", ""),
        (r"^sino\s+si\s+", ""),
        (r"^sino\s+ang\s+", ""),
        (r"^kailan\s+", ""),
        (r"^bakit\s+", ""),
        (r"^paano\s+", ""),
    ]

    topic = text
    for pattern, replacement in replacements:
        topic = re.sub(pattern, replacement, topic, flags=re.IGNORECASE).strip()

    topic = re.sub(r"\b(about|regarding|related to)\b", "", topic, flags=re.IGNORECASE).strip()
    topic = re.sub(r"\s+", " ", topic).strip(" -:;,")

    if not topic or topic.lower() == text.lower():
        topic = text

    # Avoid using the exact question wording as a document title.
    if topic.lower() == text.lower():
        title = f"{title_case_words(topic)} Summary"
    else:
        title = title_case_words(topic)

    if len(title) > 70:
        title = title[:70].rsplit(" ", 1)[0].strip()

    return title or "InknowVa Response"


def make_filename_slug(title):
    # Safe, short filename part based on the generated topic title.
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", str(title or "inknowva_response").lower()).strip("_")
    slug = re.sub(r"_+", "_", slug)
    return slug[:60].strip("_") or "inknowva_response"


def make_export_filename(export_type, title=None):
    # Dynamic filename based on the topic title, not the exact question.
    extension, _mime = EXPORT_TYPES[export_type]
    slug = make_filename_slug(title or "InknowVa Response")
    return f"{slug}_{get_export_timestamp()}.{extension}"


def get_question_for_message(message_index=None):
    # Pull the original question from the assistant message when available.
    if message_index is None:
        return ""

    messages = st.session_state.get("messages", [])

    try:
        message = messages[int(message_index)]
    except Exception:
        return ""

    if isinstance(message, dict):
        question = str(message.get("question") or "").strip()

        if question:
            return question

    # Fallback: use nearest previous user message.
    try:
        start_index = int(message_index) - 1
    except Exception:
        return ""

    for index in range(start_index, -1, -1):
        previous_message = messages[index]

        if isinstance(previous_message, dict) and previous_message.get("role") == "user":
            return str(previous_message.get("content") or "").strip()

    return ""


def normalize_answer_text(answer):
    # Normalize answer text while preserving paragraphs and list lines.
    text = str(answer or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text


def is_bullet_line(line):
    # Detect bullets from Markdown/plain text.
    return bool(re.match(r"^\s*([-*•])\s+", str(line or "")))


def is_numbered_line(line):
    # Detect numbered list items such as "1. item" or "1) item".
    return bool(re.match(r"^\s*\d+[\.\)]\s+", str(line or "")))


def strip_list_marker(line):
    # Remove bullet/number markers before applying document list styles.
    line = str(line or "").strip()
    line = re.sub(r"^\s*[-*•]\s+", "", line)
    line = re.sub(r"^\s*\d+[\.\)]\s+", "", line)
    return line.strip()


def split_answer_lines(answer):
    # Return meaningful lines while keeping blank lines as paragraph breaks.
    text = normalize_answer_text(answer)

    if not text:
        return []

    return text.split("\n")


def get_source_entries(sources):
    # Normalize and deduplicate sources for all export formats.
    entries = []

    for source in deduplicate_sources(sources or [], max_sources=None):
        item = normalize_source_item(source) or {}

        if not item:
            continue

        name = str(item.get("source") or "Unknown source").strip()
        page = str(item.get("page") or "N/A").strip()
        preview = clean_preview_text(item.get("preview"), limit=220)
        entries.append({
            "source": name,
            "page": page,
            "preview": preview,
        })

    return entries


def format_source_line(source, index=None):
    # Source line for export files.
    source = normalize_source_item(source) or {}
    prefix = f"{index}. " if index is not None else ""
    return f"{prefix}{source.get('source', 'Unknown source')} | Page: {source.get('page', 'N/A')}"


def build_export_text(answer, sources=None, title=None):
    # Shared plain-text body for text-style exports.
    # Exported files include only the topic title and answer; sources stay in the app UI.
    lines = [
        str(title or "InknowVa Response"),
        "",
        normalize_answer_text(answer) or "No answer.",
    ]

    return "\n".join(lines)


def export_txt(answer, sources=None, title=None):
    return build_export_text(answer=answer, sources=sources, title=title).encode("utf-8")


def export_md(answer, sources=None, title=None):
    lines = [
        f"# {title or 'InknowVa Response'}",
        "",
        normalize_answer_text(answer) or "No answer.",
    ]

    return "\n".join(lines).encode("utf-8")


def export_csv(answer, sources=None, title=None):
    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(["Title", "Answer"])
    writer.writerow([str(title or "InknowVa Response"), normalize_answer_text(answer)])

    return output.getvalue().encode("utf-8-sig")


def export_xlsx(answer, sources=None, title=None):
    if Workbook is None:
        return None

    file = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Response"

    sheet["A1"] = str(title or "InknowVa Response")
    sheet["A3"] = "Answer"
    sheet["A4"] = normalize_answer_text(answer) or "No answer."

    sheet.column_dimensions["A"].width = 100

    if Font is not None:
        sheet["A1"].font = Font(bold=True, size=16)
        sheet["A3"].font = Font(bold=True, size=12)

    if Alignment is not None:
        sheet["A1"].alignment = Alignment(wrap_text=True, vertical="top")
        sheet["A4"].alignment = Alignment(wrap_text=True, vertical="top")

    workbook.save(file)
    file.seek(0)
    return file


def export_xls(answer, sources=None, title=None):
    rows = [
        [str(title or "InknowVa Response")],
        [""],
        ["Answer"],
        [normalize_answer_text(answer) or "No answer."],
    ]

    html_rows = []
    for row_index, row in enumerate(rows):
        tag = "th" if row_index in [0, 2] else "td"
        cells = "".join(
            f"<{tag} style='padding:8px;border:1px solid #cccccc;vertical-align:top;text-align:left;white-space:pre-wrap;'>{html.escape(str(cell))}</{tag}>"
            for cell in row
        )
        html_rows.append(f"<tr>{cells}</tr>")

    html_doc = "".join([
        "<html><head><meta charset='utf-8'></head><body>",
        "<table style='border-collapse:collapse;font-family:Arial,sans-serif;font-size:12px;width:100%;'>",
        "".join(html_rows),
        "</table>",
        "</body></html>",
    ])

    return html_doc.encode("utf-8")



def add_docx_paragraph(doc, line):
    # Add one line to DOCX while preserving bullets and numbered items.
    raw_line = str(line or "")

    if not raw_line.strip():
        doc.add_paragraph("")
        return

    if is_bullet_line(raw_line):
        doc.add_paragraph(strip_list_marker(raw_line), style="List Bullet")
        return

    if is_numbered_line(raw_line):
        doc.add_paragraph(strip_list_marker(raw_line), style="List Number")
        return

    paragraph = doc.add_paragraph(raw_line.strip())
    paragraph_format = paragraph.paragraph_format
    paragraph_format.space_after = DocxPt(6) if DocxPt else None


def export_docx(answer, sources=None, title=None):
    if Document is None:
        return None

    file = BytesIO()
    doc = Document()

    if DocxInches is not None:
        for section in doc.sections:
            section.top_margin = DocxInches(0.75)
            section.bottom_margin = DocxInches(0.75)
            section.left_margin = DocxInches(0.75)
            section.right_margin = DocxInches(0.75)

    heading = doc.add_heading(str(title or "InknowVa Response"), level=0)

    if WD_ALIGN_PARAGRAPH is not None:
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    answer_lines = split_answer_lines(answer)

    if answer_lines:
        for line in answer_lines:
            add_docx_paragraph(doc, line)
    else:
        doc.add_paragraph("No answer.")

    doc.save(file)
    file.seek(0)
    return file



def get_pdf_styles():
    # Build reportlab styles with safe fallbacks.
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="AnswerText",
        parent=styles["BodyText"],
        fontSize=10,
        leading=14,
        spaceAfter=7,
    ))
    styles.add(ParagraphStyle(
        name="BulletText",
        parent=styles["BodyText"],
        leftIndent=16,
        firstLineIndent=-8,
        fontSize=10,
        leading=14,
        spaceAfter=5,
    ))
    return styles


def add_pdf_answer_lines(story, styles, answer):
    # Add answer lines to PDF while preserving simple list formatting.
    answer_lines = split_answer_lines(answer)

    if not answer_lines:
        story.append(Paragraph("No answer.", styles["AnswerText"]))
        return

    for line in answer_lines:
        raw_line = str(line or "")

        if not raw_line.strip():
            story.append(Spacer(1, 6))
            continue

        escaped = html.escape(strip_list_marker(raw_line) if (is_bullet_line(raw_line) or is_numbered_line(raw_line)) else raw_line.strip())

        if is_bullet_line(raw_line):
            story.append(Paragraph(f"• {escaped}", styles["BulletText"]))
        elif is_numbered_line(raw_line):
            story.append(Paragraph(escaped, styles["BulletText"]))
        else:
            story.append(Paragraph(escaped, styles["AnswerText"]))


def export_pdf(answer, sources=None, title=None):
    if None in [getSampleStyleSheet, Paragraph, SimpleDocTemplate, Spacer, ParagraphStyle, colors]:
        return None

    file = BytesIO()
    pdf = SimpleDocTemplate(
        file,
        pagesize=letter,
        rightMargin=42,
        leftMargin=42,
        topMargin=42,
        bottomMargin=42,
    )
    styles = get_pdf_styles()
    story = [
        Paragraph(html.escape(str(title or "InknowVa Response")), styles["Title"]),
        Spacer(1, 10),
    ]

    add_pdf_answer_lines(story, styles, answer)

    pdf.build(story)
    file.seek(0)
    return file



def split_text_for_slides(text, chunk_size=650):
    text = normalize_answer_text(text)

    if not text:
        return ["No answer."]

    chunks = []
    while text:
        chunk = text[:chunk_size]
        cut = max(chunk.rfind(". "), chunk.rfind("\n"), chunk.rfind(" "))

        if cut > 180:
            chunk = text[:cut + 1]

        chunks.append(chunk.strip())
        text = text[len(chunk):].strip()

    return chunks


def export_pptx(answer, sources=None, title=None):
    if Presentation is None:
        return None

    file = BytesIO()
    presentation = Presentation()

    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = str(title or "InknowVa Response")
    slide.placeholders[1].text = "Generated by InknowVa"

    content_layout = presentation.slide_layouts[1]
    for index, chunk in enumerate(split_text_for_slides(answer), start=1):
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = "Answer" if index == 1 else f"Answer {index}"
        body = slide.placeholders[1]
        body.text = chunk

        if PptxPt is not None:
            for paragraph in body.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = PptxPt(16)

    presentation.save(file)
    file.seek(0)
    return file


def export_ppt(answer, sources=None, title=None):
    rows = [
        str(title or "InknowVa Response"),
        "",
        normalize_answer_text(answer) or "No answer.",
    ]

    # Legacy .ppt binary generation is not supported directly.
    # This HTML-based .ppt opens in PowerPoint/Office with a compatibility warning.
    body = "<br>".join(html.escape(row) for row in rows)
    return f"<html><body>{body}</body></html>".encode("utf-8")



def get_export_file(answer, sources, export_type, title=None):
    # Return file data, filename, and MIME.
    exporters = {
        "PDF": export_pdf,
        "DOCX": export_docx,
        "TXT": export_txt,
        "Markdown": export_md,
        "CSV": export_csv,
        "XLSX": export_xlsx,
        "XLS": export_xls,
        "PPTX": export_pptx,
        "PPT": export_ppt,
    }

    if export_type not in EXPORT_TYPES:
        return None, None, None

    exporter = exporters.get(export_type)

    if not exporter:
        return None, None, None

    data = exporter(answer, sources, title=title)

    if data is None:
        return None, None, None

    filename = make_export_filename(export_type, title=title)
    _extension, mime = EXPORT_TYPES[export_type]

    return data, filename, mime


def set_regenerate_index(index):
    # Mark assistant message for full RAG regeneration in chat_ui.py.
    # Hide answer-only UI immediately on the click rerun.
    append_action_debug("ACTION UI REGENERATE CLICKED", [f"index: {index}"])
    st.session_state.regenerate_index = index
    st.session_state.regenerate_requested_at = time.time()
    st.session_state.hide_actions = True
    st.session_state.is_generating = True
    st.session_state.pop("open_export_panel_key", None)


def toggle_export_panel(unique_key):
    # Open/close the inline export panel for one answer.
    current_key = st.session_state.get("open_export_panel_key")

    if str(current_key) == str(unique_key):
        st.session_state.pop("open_export_panel_key", None)
        clear_export_runtime_state()
        return

    clear_export_runtime_state()
    st.session_state.open_export_panel_key = str(unique_key)


def clear_export_runtime_state():
    # Clear export request/result state when switching panels.
    st.session_state.pop(PENDING_EXPORT_REQUEST_KEY, None)
    st.session_state.pop(READY_EXPORT_FILE_KEY, None)


def queue_export_request(unique_key, export_type, title=None):
    # Store which export the user wants. Generation happens after a loading state renders.
    st.session_state[PENDING_EXPORT_REQUEST_KEY] = {
        "unique_key": str(unique_key),
        "export_type": str(export_type),
        "title": title,
    }
    st.session_state.pop(READY_EXPORT_FILE_KEY, None)


def get_pending_export_request(unique_key):
    # Return pending export request for this answer only.
    request = st.session_state.get(PENDING_EXPORT_REQUEST_KEY)

    if not isinstance(request, dict):
        return None

    if str(request.get("unique_key")) != str(unique_key):
        return None

    return request


def get_ready_export_file(unique_key):
    # Return prepared export bytes for this answer only.
    prepared = st.session_state.get(READY_EXPORT_FILE_KEY)

    if not isinstance(prepared, dict):
        return None

    if str(prepared.get("unique_key")) != str(unique_key):
        return None

    return prepared


def display_export_loading(export_type):
    # Show loading before generating files like PDF/DOCX/XLSX/PPTX.
    st.markdown(
        f"""
<div class="export-loading-row">
    <div class="thinking-dots"><span></span><span></span><span></span></div>
    <span>Preparing {safe_html(export_type)}...</span>
</div>
""",
        unsafe_allow_html=True,
    )

    if EXPORT_FRONTEND_FLUSH_SECONDS > 0:
        time.sleep(EXPORT_FRONTEND_FLUSH_SECONDS)


def prepare_pending_export(answer, sources, unique_key, title=None):
    # Generate only the requested export type, not every format at once.
    request = get_pending_export_request(unique_key)

    if not request:
        return

    export_type = request.get("export_type")
    display_export_loading(export_type)

    data, filename, mime = get_export_file(
        answer=answer,
        sources=sources,
        export_type=export_type,
        title=request.get("title") or title,
    )

    if data is not None:
        st.session_state[READY_EXPORT_FILE_KEY] = {
            "unique_key": str(unique_key),
            "export_type": export_type,
            "data": data,
            "filename": filename,
            "mime": mime,
        }

    st.session_state.pop(PENDING_EXPORT_REQUEST_KEY, None)
    st.rerun()


def display_ready_export(unique_key):
    # Show the final download button after the file is prepared.
    prepared = get_ready_export_file(unique_key)

    if not prepared:
        return

    export_type = prepared.get("export_type") or "file"

    st.download_button(
        label=f"Download {export_type}",
        data=prepared.get("data"),
        file_name=prepared.get("filename"),
        mime=prepared.get("mime"),
        key=f"download_ready_{str(export_type).lower()}_{unique_key}",
        use_container_width=True,
    )


def display_export_buttons(answer, sources, unique_key, title=None):
    # Export panel is fast: it renders format buttons first.
    # The actual file is generated only after a specific format is clicked.
    options = [
        ("PDF", "PDF"),
        ("DOCX", "DOCX"),
        ("TXT", "TXT"),
        ("Markdown", "MD"),
        ("XLSX", "XLSX"),
        ("XLS", "XLS"),
        ("CSV", "CSV"),
        ("PPTX", "PPTX"),
        ("PPT", "PPT"),
    ]

    columns = st.columns(len(options))

    for column, option in zip(columns, options):
        export_type, label = option

        with column:
            st.button(
                label,
                key=f"prepare_export_{label.lower()}_{unique_key}",
                use_container_width=True,
                on_click=queue_export_request,
                args=(unique_key, export_type, title),
            )

    prepare_pending_export(answer, sources, unique_key, title=title)
    display_ready_export(unique_key)


def prepare_sources_for_drawer(sources):
    # Clean and deduplicate source files.
    return deduplicate_sources(sources or [], max_sources=None)


def get_optional_source_url(source):
    # Use web/static URL when source metadata already has one.
    if not isinstance(source, dict):
        return ""

    metadata = source.get("metadata") or {}

    for key in ["url", "link", "file_url", "source_url", "path_url"]:
        value = source.get(key) or metadata.get(key)

        if value and str(value).strip().startswith(("http://", "https://", "/")):
            return str(value).strip()

    return ""


def build_source_card_html(index, source):
    # Build one source drawer card.
    source = normalize_source_item(source)

    if not source:
        return ""

    name = source.get("source", "Unknown source")
    page = source.get("page", "N/A")
    preview = clean_preview_text(source.get("preview"), limit=SOURCE_PREVIEW_LIMIT)
    link = make_static_file_link(source.get("source_path"), page=page) or get_optional_source_url(source)

    safe_name = safe_html(name)
    safe_page = safe_html(page)
    safe_preview = safe_html(preview)

    if link:
        safe_link = safe_html(link)
        title_html = f'<a class="right-source-name" href="{safe_link}" target="_blank" rel="noopener noreferrer">{safe_name}</a>'
        preview_html = f'<a class="right-source-preview" href="{safe_link}" target="_blank" rel="noopener noreferrer">{safe_preview}</a>'
    else:
        title_html = f'<span class="right-source-name">{safe_name}</span>'
        preview_html = f'<div class="right-source-preview right-source-preview-disabled">{safe_preview}</div>'

    return "\n".join([
        '<div class="right-source-card">',
        '<div class="right-source-card-header">',
        f'<span class="right-source-number">{index}</span>',
        title_html,
        '</div>',
        f'<div class="right-source-page">Page: {safe_page}</div>',
        preview_html,
        '</div>',
    ])


def render_source_cards_html(sources):
    # Render drawer source card HTML.
    cards = []

    for index, source in enumerate(prepare_sources_for_drawer(sources), start=1):
        card_html = build_source_card_html(index, source)

        if card_html:
            cards.append(card_html)

    if not cards:
        return '<div class="right-source-empty">No sources available for this answer.</div>'

    return "\n".join(cards)


def show_sources_in_sidebar(sources):
    # Store current answer sources for the sidebar Sources block.
    st.session_state[SIDEBAR_SOURCES_KEY] = prepare_sources_for_drawer(sources)


def display_sources_button(sources, unique_key):
    # Send current answer sources to the left sidebar instead of opening a drawer.
    st.button(
        "▤",
        key=f"sources_{unique_key}",
        help="Show sources in sidebar",
        disabled=not bool(sources),
        on_click=show_sources_in_sidebar,
        args=(sources,),
    )


def display_actions(answer, sources, message_index=None):
    # Copy, regenerate, and export buttons below one assistant answer.
    # Sources are already shown in the left sidebar, so no source button is needed here.
    if st.session_state.get("hide_actions", False) or st.session_state.get("is_generating", False):
        return

    unique_key = str(0 if message_index is None else message_index)
    export_title = make_response_title(get_question_for_message(message_index))

    # Do not update session_state while rendering action rows.
    # chat_ui.py already stores latest sources after generation/history load.

    with st.container(key=f"action_area_{unique_key}"):
        col1, col2, col3, col4 = st.columns([1, 1, 1, 12], gap="small")

        with col1:
            copy_button(answer, key=f"copy_{unique_key}")

        with col2:
            st.button(
                "↻",
                key=f"regen_{unique_key}",
                help="Regenerate answer",
                on_click=set_regenerate_index,
                args=(message_index,),
            )

        with col3:
            st.button(
                "⇩",
                key=f"export_toggle_{unique_key}",
                help="Download response",
                on_click=toggle_export_panel,
                args=(unique_key,),
            )

        if st.session_state.get("open_export_panel_key") == unique_key:
            with col4:
                with st.container(key=f"export_inline_panel_{unique_key}"):
                    display_export_buttons(answer, sources, unique_key, title=export_title)
